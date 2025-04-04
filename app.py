import os
import sys
from flask import Flask, request, jsonify, render_template, send_file, session
from werkzeug.utils import secure_filename
import pytesseract
from PIL import Image
import pdf2image
from pdf2image.exceptions import PDFInfoNotInstalledError
import docx
import magic
import logging
from pathlib import Path
import subprocess
import tempfile
from datetime import datetime, timedelta
import hashlib
import pandas as pd
from pptx import Presentation
import ebooklib
from ebooklib import epub
from bs4 import BeautifulSoup

# Configure logging
logging.basicConfig(level=logging.DEBUG,
                   format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

app = Flask(__name__, static_folder='static')
app.config['MAX_CONTENT_LENGTH'] = 20 * 1024 * 1024  # 20MB max file size
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['TEMP_FOLDER'] = 'temp'
app.secret_key = os.urandom(24)

# Configure paths - Update Poppler path to match user's installation
POPPLER_PATH = r'D:\cursor\Release-24.08.0-0\poppler-24.08.0\Library\bin'
os.environ['POPPLER_PATH'] = POPPLER_PATH

# Ensure upload and temp folders exist
UPLOAD_FOLDER = Path(app.config['UPLOAD_FOLDER'])
TEMP_FOLDER = Path(app.config['TEMP_FOLDER'])
UPLOAD_FOLDER.mkdir(exist_ok=True)
TEMP_FOLDER.mkdir(exist_ok=True)

# Configure Tesseract path - including D: drive locations
TESSERACT_PATHS = [
    r'D:\Tesseract-OCR\tesseract.exe',
    r'D:\Program Files\Tesseract-OCR\tesseract.exe',
    r'C:\Program Files\Tesseract-OCR\tesseract.exe',
    r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
]

ALLOWED_EXTENSIONS = {
    # Images
    'png', 'jpg', 'jpeg', 'bmp', 'tiff', 'gif',
    # Documents
    'pdf', 'docx', 'doc', 'txt', 'rtf',
    # Spreadsheets
    'xlsx', 'xls', 'csv',
    # Presentations
    'pptx', 'ppt',
    # Other
    'odt', 'ods', 'odp', 'epub', 'mobi'
}

def verify_poppler():
    """Verify Poppler installation"""
    try:
        if not os.path.exists(POPPLER_PATH):
            logger.error(f"Poppler directory not found at: {POPPLER_PATH}")
            return False

        required_files = ['pdfinfo.exe', 'pdftoppm.exe']
        missing_files = []
        for file in required_files:
            file_path = os.path.join(POPPLER_PATH, file)
            if not os.path.exists(file_path):
                missing_files.append(file)

        if missing_files:
            logger.error(f"Missing Poppler files: {', '.join(missing_files)}")
            return False

        test_pdf = os.path.join(POPPLER_PATH, 'pdfinfo.exe')
        if os.path.exists(test_pdf):
            logger.info(f"Poppler successfully verified at: {POPPLER_PATH}")
            return True
        else:
            logger.error("Could not verify Poppler installation")
            return False
    except Exception as e:
        logger.error(f"Error verifying Poppler: {str(e)}")
        return False

def find_tesseract():
    """Find Tesseract executable"""
    for path in TESSERACT_PATHS:
        if os.path.exists(path):
            logger.info(f"Found Tesseract at: {path}")
            return path
            
    try:
        if os.name == 'nt':  # Windows
            result = subprocess.run(['where', 'tesseract'], capture_output=True, text=True)
            if result.returncode == 0:
                path = result.stdout.strip()
                logger.info(f"Found Tesseract in PATH: {path}")
                return path
    except Exception as e:
        logger.error(f"Error finding tesseract in PATH: {e}")
    
    logger.error("Tesseract not found. Please verify the installation path.")
    return None

def setup_tesseract():
    """Setup Tesseract configuration"""
    tesseract_path = find_tesseract()
    if tesseract_path:
        pytesseract.pytesseract.tesseract_cmd = tesseract_path
        try:
            version = pytesseract.get_tesseract_version()
            logger.info(f"Tesseract version: {version}")
            return True
        except Exception as e:
            logger.error(f"Error verifying Tesseract installation: {e}")
            return False
    return False

def check_dependencies():
    """Check if all required dependencies are properly installed"""
    errors = []
    
    try:
        version = pytesseract.get_tesseract_version()
        logger.info(f"Tesseract version: {version}")
    except Exception as e:
        error_msg = f"Tesseract error: {str(e)}"
        logger.error(error_msg)
        errors.append(error_msg)

    if not verify_poppler():
        error_msg = "Poppler not properly configured. Please check installation instructions."
        logger.error(error_msg)
        errors.append(error_msg)

    if errors:
        return False, errors
    return True, []

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def get_file_hash(file_path):
    """Generate a hash of the file contents"""
    hasher = hashlib.md5()
    with open(file_path, 'rb') as f:
        buf = f.read(65536)  # Read in 64kb chunks
        while len(buf) > 0:
            hasher.update(buf)
            buf = f.read(65536)
    return hasher.hexdigest()

def save_result_to_temp(text, file_hash):
    """Save extracted text to a temporary file"""
    temp_file = TEMP_FOLDER / f"{file_hash}.txt"
    with open(temp_file, 'w', encoding='utf-8') as f:
        f.write(text)
    return temp_file

def extract_text_from_image(image_path):
    try:
        logger.debug(f"Opening image: {image_path}")
        with Image.open(image_path) as img:
            # Optimize image size if too large
            max_size = 2000
            if max(img.size) > max_size:
                ratio = max_size / max(img.size)
                new_size = tuple(int(dim * ratio) for dim in img.size)
                img = img.resize(new_size, Image.Resampling.LANCZOS)
            
            # Convert to grayscale and optimize for OCR
            img = img.convert('L')
            
            # Use a faster thresholding method
            img = img.point(lambda x: 0 if x < 128 else 255, '1')
            
            logger.debug("Performing OCR on image")
            text = pytesseract.image_to_string(img, config='--oem 3 --psm 6')
            logger.debug(f"OCR completed, extracted {len(text)} characters")
            return text
    except Exception as e:
        logger.error(f"Error in image extraction: {str(e)}")
        raise

def extract_text_from_pdf(pdf_path):
    """Extract text from PDF using Poppler"""
    try:
        if not verify_poppler():
            raise Exception("Poppler is not properly configured")

        logger.info(f"Converting PDF using Poppler at: {POPPLER_PATH}")
        images = pdf2image.convert_from_path(
            pdf_path,
            poppler_path=POPPLER_PATH,
            dpi=200,  # Reduced DPI for faster processing
            thread_count=4  # Enable multi-threading
        )
        
        logger.info(f"Successfully converted {len(images)} pages")
        text = ""
        for i, image in enumerate(images, 1):
            logger.info(f"Processing page {i}/{len(images)}")
            # Optimize image size if too large
            max_size = 2000
            if max(image.size) > max_size:
                ratio = max_size / max(image.size)
                new_size = tuple(int(dim * ratio) for dim in image.size)
                image = image.resize(new_size, Image.Resampling.LANCZOS)
            
            # Convert to grayscale and optimize for OCR
            image = image.convert('L')
            image = image.point(lambda x: 0 if x < 128 else 255, '1')
            
            text += f"[Page {i}]\n{pytesseract.image_to_string(image, config='--oem 3 --psm 6')}\n\n"
        return text
    except PDFInfoNotInstalledError:
        logger.error("Poppler is not properly installed")
        raise Exception("PDF processing failed: Poppler not properly installed")
    except Exception as e:
        logger.error(f"Error processing PDF: {str(e)}")
        raise

def extract_text_from_docx(docx_path):
    try:
        logger.debug(f"Opening DOCX file: {docx_path}")
        doc = docx.Document(docx_path)
        
        # Extract text with minimal formatting
        text = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():  # Skip empty paragraphs
                text.append(paragraph.text)
        
        formatted_text = "\n".join(text)
        logger.debug(f"Extracted {len(formatted_text)} characters from DOCX")
        return formatted_text
    except Exception as e:
        logger.error(f"Error in DOCX extraction: {str(e)}")
        raise

def extract_text_from_text_file(file_path):
    """Extract text from plain text files"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        # Try different encodings if UTF-8 fails
        encodings = ['latin1', 'cp1252', 'iso-8859-1']
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        raise Exception("Could not decode text file with any supported encoding")

def extract_text_from_excel(file_path):
    """Extract text from Excel files"""
    try:
        # Read only the first sheet by default
        df = pd.read_excel(file_path, sheet_name=0, nrows=1000)  # Limit rows for faster processing
        return df.to_string()
    except ImportError:
        logger.error("pandas not installed")
        raise Exception("Excel processing requires pandas to be installed")
    except Exception as e:
        logger.error(f"Error extracting text from Excel: {str(e)}")
        raise

def extract_text_from_presentation(file_path):
    """Extract text from PowerPoint files"""
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            if slide_text:  # Only add non-empty slides
                text.append(f"\n[Slide {slide.slide_id}]\n" + "\n".join(slide_text))
        return "\n".join(text)
    except ImportError:
        logger.error("python-pptx not installed")
        raise Exception("Presentation processing requires python-pptx to be installed")
    except Exception as e:
        logger.error(f"Error extracting text from presentation: {str(e)}")
        raise

def extract_text_from_epub(file_path):
    """Extract text from EPUB files"""
    try:
        book = epub.read_epub(file_path)
        text = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                text.append(soup.get_text())
        return "\n".join(text)
    except ImportError:
        logger.error("ebooklib or beautifulsoup4 not installed")
        raise Exception("EPUB processing requires ebooklib and beautifulsoup4 to be installed")
    except Exception as e:
        logger.error(f"Error extracting text from EPUB: {str(e)}")
        raise

@app.before_request
def before_request():
    session.permanent = True
    app.permanent_session_lifetime = timedelta(days=1)

@app.errorhandler(413)
def request_entity_too_large(error):
    return jsonify({'error': 'File too large'}), 413

@app.errorhandler(500)
def internal_server_error(error):
    return jsonify({'error': 'Internal server error'}), 500

def cleanup_temp_files():
    """Clean up old temporary files"""
    try:
        current_time = datetime.now()
        for file in TEMP_FOLDER.glob('*.txt'):
            file_age = datetime.fromtimestamp(file.stat().st_mtime)
            if (current_time - file_age).days > 1:  # Delete files older than 1 day
                file.unlink()
                logger.info(f"Cleaned up old temp file: {file}")
    except Exception as e:
        logger.error(f"Error cleaning up temp files: {str(e)}")

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    logger.debug("Received file upload request")
    
    # Check dependencies
    deps_ok, errors = check_dependencies()
    if not deps_ok:
        error_msg = "Server configuration error: " + "; ".join(errors)
        logger.error(error_msg)
        return jsonify({'error': error_msg}), 500

    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    
    if not allowed_file(file.filename):
        return jsonify({'error': 'File type not allowed'}), 400
    
    try:
        # Check file size before processing
        if request.content_length > app.config['MAX_CONTENT_LENGTH']:
            return jsonify({'error': 'File too large'}), 413

        filename = secure_filename(file.filename)
        filepath = UPLOAD_FOLDER / filename
        file.save(str(filepath))
        
        # Generate file hash for caching
        file_hash = get_file_hash(filepath)
        
        # Check if we have cached results
        cached_result = TEMP_FOLDER / f"{file_hash}.txt"
        mime = magic.Magic(mime=True)
        file_type = None
        
        if cached_result.exists():
            logger.info("Using cached result")
            with open(cached_result, 'r', encoding='utf-8') as f:
                text = f.read()
            extraction_time = 0  # No time for cached results
        else:
            file_type = mime.from_file(str(filepath))
            logger.debug(f"Processing file: {filename} (type: {file_type})")
            
            # Start timing the extraction
            start_time = datetime.now()
            
            # Determine file type and extract text accordingly
            if file_type.startswith('image/'):
                text = extract_text_from_image(str(filepath))
            elif file_type == 'application/pdf':
                text = extract_text_from_pdf(str(filepath))
            elif file_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                text = extract_text_from_docx(str(filepath))
            elif file_type == 'text/plain':
                text = extract_text_from_text_file(str(filepath))
            elif file_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'application/vnd.ms-excel']:
                text = extract_text_from_excel(str(filepath))
            elif file_type in [
                'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                'application/vnd.ms-powerpoint',
                'application/powerpoint',
                'application/mspowerpoint',
                'application/x-mspowerpoint'
            ] or filename.lower().endswith('.pptx'):
                logger.info(f"Processing PPTX file: {filename}")
                text = extract_text_from_presentation(str(filepath))
            elif file_type == 'application/epub+zip':
                text = extract_text_from_epub(str(filepath))
            else:
                logger.error(f"Unsupported file type: {file_type} for file: {filename}")
                return jsonify({'error': f'Unsupported file type: {file_type}'}), 400
            
            # Calculate extraction time
            extraction_time = (datetime.now() - start_time).total_seconds()
            
            # Cache the results
            save_result_to_temp(text, file_hash)
        
        if not text or text.strip() == '':
            return jsonify({'error': 'No text could be extracted from the file'}), 400
            
        return jsonify({
            'text': text,
            'metadata': {
                'filename': filename,
                'processed_at': datetime.now().isoformat(),
                'file_size': os.path.getsize(filepath),
                'file_type': file_type,
                'extraction_time': round(extraction_time, 2)  # Round to 2 decimal places
            }
        })
    
    except Exception as e:
        logger.error(f"Error processing file: {str(e)}", exc_info=True)
        return jsonify({'error': str(e)}), 500
    
    finally:
        if 'filepath' in locals() and filepath.exists():
            try:
                filepath.unlink()
            except Exception as e:
                logger.error(f"Error cleaning up file: {str(e)}")
        # Clean up old temporary files
        cleanup_temp_files()

@app.route('/download/<file_hash>')
def download_result(file_hash):
    """Download extracted text as a file"""
    try:
        file_path = TEMP_FOLDER / f"{file_hash}.txt"
        if not file_path.exists():
            return jsonify({'error': 'File not found'}), 404
        
        return send_file(
            file_path,
            as_attachment=True,
            download_name=f"extracted_text_{file_hash[:8]}.txt",
            mimetype='text/plain'
        )
    except Exception as e:
        logger.error(f"Error downloading file: {str(e)}")
        return jsonify({'error': 'Failed to download file'}), 500

if __name__ == '__main__':
    logger.info("Starting Text Extractor application")
    logger.info(f"Python version: {sys.version}")
    logger.info(f"Operating system: {os.name}")
    
    if not setup_tesseract():
        logger.error("Failed to configure Tesseract. Please check installation.")
    
    if not verify_poppler():
        logger.error("Failed to configure Poppler. Please check installation.")
    
    deps_ok, errors = check_dependencies()
    if not deps_ok:
        logger.error("Dependency check failed:")
        for error in errors:
            logger.error(f"  - {error}")
    
    app.run(debug=True) 